"""Gera a apresentacao em PPTX para o Grupo 16 — AVL Aumentada."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Paleta de cores
# ---------------------------------------------------------------------------
DARK   = RGBColor(0x1A, 0x1A, 0x2E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
AVL_C  = RGBColor(0x17, 0x6B, 0x87)
BST_C  = RGBColor(0xC4, 0x45, 0x36)
ACCENT = RGBColor(0xF5, 0xA6, 0x23)
GRAY   = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

GRAFICOS = Path(__file__).parent.parent / "graficos"


def rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def set_bg(slide, color: RGBColor) -> None:
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    text: str,
    left: float, top: float, width: float, height: float,
    font_size: int = 20,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
    wrap: bool = True,
) -> None:
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_para(tf, text: str, font_size: int = 18, bold: bool = False,
             color: RGBColor = WHITE, align=PP_ALIGN.LEFT, indent: int = 0) -> None:
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_content_box(
    slide,
    left: float, top: float, width: float, height: float,
    font_size: int = 18,
) -> object:
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    # Limpa paragrafo inicial
    tf.paragraphs[0].text = ""
    return tf


def add_image(slide, path: Path, left: float, top: float, width: float) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
    else:
        add_textbox(slide, f"[imagem nao encontrada: {path.name}]",
                    left, top, width, 3, font_size=12, color=rgb(255, 80, 80))


def section_bar(slide, title: str, color: RGBColor = AVL_C) -> None:
    """Barra colorida no topo com titulo da secao."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.33), Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "  " + title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = WHITE


# ---------------------------------------------------------------------------
# SLIDES
# ---------------------------------------------------------------------------

def slide_capa(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)

    # Linha decorativa AVL
    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.8), Inches(13.33), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AVL_C
    bar.line.fill.background()

    add_textbox(slide, "AVL Aumentada", 0.5, 1.0, 12, 1.4, font_size=46, bold=True,
                color=AVL_C, align=PP_ALIGN.LEFT)
    add_textbox(slide, "Rank · Select · Range Aggregation em O(log n)",
                0.5, 2.2, 12, 0.8, font_size=22, color=GRAY, align=PP_ALIGN.LEFT)

    add_textbox(slide, "Estruturas de Dados — Grupo 16", 0.5, 3.3, 12, 0.6,
                font_size=18, color=WHITE)
    add_textbox(slide, "Universidade Federal do Pampa · 2026", 0.5, 3.9, 12, 0.6,
                font_size=16, color=GRAY)

    add_textbox(slide, "Parte 1  Código e implementação", 0.5, 5.2, 6, 0.6,
                font_size=17, color=AVL_C, bold=True)
    add_textbox(slide, "Parte 2  Experimentos e gráficos", 6.8, 5.2, 6, 0.6,
                font_size=17, color=ACCENT, bold=True)


def slide_node_structure(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Estrutura do Nó — AVL Aumentada", AVL_C)

    tf = add_content_box(slide, 0.5, 1.3, 5.8, 5.8, font_size=16)
    items = [
        ("Campos do nó:", True, 18, AVL_C),
        ("key   — chave inteira", False, 16, WHITE),
        ("left, right  — filhos", False, 16, WHITE),
        ("height  — altura da subárvore", False, 16, WHITE),
        ("size  — nº de nós na subárvore", False, 16, WHITE),
        ("subtree_sum  — soma das chaves", False, 16, WHITE),
        ("", False, 10, WHITE),
        ("Invariantes mantidos:", True, 18, ACCENT),
        ("BST: esq < raiz < dir", False, 16, WHITE),
        ("AVL: |h_esq − h_dir| ≤ 1", False, 16, WHITE),
        ("size = size(esq) + size(dir) + 1", False, 16, WHITE),
        ("sum = sum(esq) + sum(dir) + key", False, 16, WHITE),
    ]
    for text, bold, size, color in items:
        add_para(tf, text, font_size=size, bold=bold, color=color)

    # Caixa de código
    code = (
        "@dataclass\n"
        "class Node:\n"
        "    key: int\n"
        "    left:  Node | None = None\n"
        "    right: Node | None = None\n"
        "    height: int = 1\n"
        "    size:   int = 1\n"
        "    subtree_sum: int = 0\n\n"
        "    def __post_init__(self):\n"
        "        self.subtree_sum = self.key"
    )
    add_textbox(slide, code, 6.7, 1.3, 6.3, 5.5, font_size=14,
                color=rgb(0xA8, 0xD8, 0xEA))


def slide_rotations(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Rotações e Invariantes Aumentados", AVL_C)

    tf = add_content_box(slide, 0.5, 1.3, 6.0, 5.8)
    bullets = [
        ("Por que a ordem importa:", True, 18, ACCENT),
        ("Na rotate_right(A):", True, 16, WHITE),
        ("  1. update(A)  — recalcula com filhos já corretos", False, 15, GRAY),
        ("  2. update(B)  — usa size(A) já correto", False, 15, GRAY),
        ("", False, 10, WHITE),
        ("Se invertida (B antes de A):", True, 16, BST_C),
        ("  B.size conta A duas vezes!", False, 15, BST_C),
        ("  size(B) = size(A_errado) + size(C) + 1", False, 15, BST_C),
        ("", False, 10, WHITE),
        ("Rotações duplas (LR / RL):", True, 16, WHITE),
        ("  = duas rotações simples em sequência", False, 15, GRAY),
        ("  ordem correta preservada em cada etapa", False, 15, GRAY),
        ("", False, 10, WHITE),
        ("Resultado:", True, 16, WHITE),
        ("  size e subtree_sum sempre consistentes", False, 15, rgb(0x90, 0xEE, 0x90)),
        ("  após qualquer rebalanceamento", False, 15, rgb(0x90, 0xEE, 0x90)),
    ]
    for text, bold, size, color in bullets:
        add_para(tf, text, font_size=size, bold=bold, color=color)

    code = (
        "def rotate_right(A: Node) -> Node:\n"
        "    B = A.left          # novo pai\n"
        "    A.left = B.right    # realoca filho\n"
        "    update(A)           # 1º: A (folha)\n"
        "    B.right = A\n"
        "    update(B)           # 2º: B (usa A)\n"
        "    return B"
    )
    add_textbox(slide, code, 6.8, 1.5, 6.2, 3.5, font_size=15,
                color=rgb(0xA8, 0xD8, 0xEA))

    add_textbox(slide, "update(node) recalcula height, size e subtree_sum\na partir dos filhos — O(1) por chamada",
                6.8, 5.2, 6.2, 1.8, font_size=14, color=GRAY)


def slide_augmented_ops(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Operações Aumentadas — O(log n)", AVL_C)

    ops = [
        ("rank(k)", AVL_C,
         "Conta nós menores que k.\nPercorre da raiz:\n"
         "  k > raiz → rank = size(esq)+1 + rank(dir,k)\n"
         "  k < raiz → rank = rank(esq, k)\n"
         "  k = raiz → rank = size(esq)"),
        ("select(i)", ACCENT,
         "Devolve a i-ésima chave em ordem.\n"
         "  r = size(esq) + 1\n"
         "  i < r → select(esq, i)\n"
         "  i > r → select(dir, i−r)\n"
         "  i = r → raiz.key"),
        ("range_agg(a,b)", rgb(0x4C, 0x95, 0x6C),
         "Soma as chaves no intervalo [a,b].\nDesce recursivamente:\n"
         "  se [a,b] cobre subárvore → subtree_sum\n"
         "  senão combina esq + raiz + dir\n"
         "  (sem varredura linear)"),
    ]
    for idx, (name, color, desc) in enumerate(ops):
        left = 0.4 + idx * 4.3
        # Caixa título
        box = slide.shapes.add_shape(1, Inches(left), Inches(1.3), Inches(4.0), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf2 = box.text_frame
        p = tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = name
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = WHITE

        add_textbox(slide, desc, left, 1.95, 4.1, 4.5, font_size=15, color=WHITE)

    add_textbox(slide,
                "BST ingênua: rank/select/range_agg percorrem toda a árvore → O(n)\n"
                "AVL aumentada: usa size e subtree_sum armazenados → O(log n)",
                0.4, 6.3, 12.5, 0.9, font_size=15, color=GRAY)


def slide_tests(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Testes e Validação", AVL_C)

    tf = add_content_box(slide, 0.5, 1.3, 12.3, 5.8)
    items = [
        ("Suite de testes unitários (unittest):", True, 19, ACCENT),
        ("  test_avl_aumentada.py  — invariantes AVL, size, subtree_sum após cada op", False, 16, WHITE),
        ("  test_bst_ingenua.py    — rank/select/range_agg na BST baseline", False, 16, WHITE),
        ("  test_executar_trace.py — pipeline completo de leitura de trace", False, 16, WHITE),
        ("", False, 10, WHITE),
        ("Validação contínua durante experimentos:", True, 19, ACCENT),
        ("  --validate-every N  chama validate_invariants() a cada N operações", False, 16, WHITE),
        ("  verifica: ordem BST, balanço AVL, size, subtree_sum em toda a árvore", False, 16, WHITE),
        ("", False, 10, WHITE),
        ("Cenários de estresse:", True, 19, ACCENT),
        ("  inserção ordenada → BST degenera em lista (h = n)", False, 16, WHITE),
        ("  inserção embaralhada + remoções → AVL mantém h ≤ 1,44 log₂ n", False, 16, WHITE),
        ("", False, 10, WHITE),
        ("Rastreabilidade:", True, 19, ACCENT),
        ("  trace .trace → execute_trace() → saída .out comparável com .expected", False, 16, WHITE),
    ]
    for text, bold, size, color in items:
        add_para(tf, text, font_size=size, bold=bold, color=color)


def slide_transition(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)

    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.1), Inches(13.33), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    add_textbox(slide, "Parte 2", 0.5, 1.5, 12, 1.0, font_size=36, bold=True,
                color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Experimentos e Análise dos Gráficos",
                0.5, 2.4, 12, 0.8, font_size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Grupo 16 · dataset: osm_cellids_800M · mix=35:30:35 · theta=0.6 · seed=16",
                0.5, 4.0, 12, 0.6, font_size=16, color=GRAY, align=PP_ALIGN.CENTER)


def slide_height(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Crescimento da Altura", ACCENT)

    img = GRAFICOS / "escala_altura.png"
    add_image(slide, img, 0.3, 1.2, 7.5)

    tf = add_content_box(slide, 8.1, 1.3, 5.0, 5.8)
    pts = [
        ("O que o gráfico mostra:", True, 17, ACCENT),
        ("Altura final após n inserções\nembaralhas (shuffle).", False, 15, WHITE),
        ("", False, 10, WHITE),
        ("AVL (azul):", True, 16, AVL_C),
        ("cresce como ≈ 1,44 log₂ n", False, 15, WHITE),
        ("n=1M → h ≈ 29", False, 15, GRAY),
        ("", False, 10, WHITE),
        ("BST (vermelho):", True, 16, BST_C),
        ("esperado ≈ 2,5 log₂ n (aleatório)", False, 15, WHITE),
        ("n=1M → h ≈ 50", False, 15, GRAY),
        ("", False, 10, WHITE),
        ("Conclusão:", True, 16, WHITE),
        ("Rebalanceamento AVL garante\náltura logarítmica em qualquer\nsequência de inserções.", False, 15, rgb(0x90, 0xEE, 0x90)),
    ]
    for text, bold, size, color in pts:
        add_para(tf, text, font_size=size, bold=bold, color=color)


def slide_basic_ops(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Escalabilidade das Operações Básicas (teoria × prática)", ACCENT)

    img = GRAFICOS / "escala_operacoes_basicas.png"
    add_image(slide, img, 0.2, 1.15, 8.8)

    tf = add_content_box(slide, 9.2, 1.3, 3.9, 5.8)
    pts = [
        ("Três painéis: insert, delete, search", True, 15, ACCENT),
        ("Curva cinza: O(log n) ancorada\nem n=100 na AVL.", False, 14, WHITE),
        ("", False, 10, WHITE),
        ("Inserção:", True, 15, AVL_C),
        ("AVL ~10× mais lenta que BST\n(overhead de rotações)", False, 14, GRAY),
        ("", False, 8, WHITE),
        ("Busca:", True, 15, AVL_C),
        ("Latências próximas;\nAVL leva vantagem em n grande\n(árvore mais baixa)", False, 14, GRAY),
        ("", False, 8, WHITE),
        ("Ambas seguem O(log n):", True, 15, rgb(0x90, 0xEE, 0x90)),
        ("crescimento sub-linear\nconfirmado empiricamente", False, 14, GRAY),
    ]
    for text, bold, size, color in pts:
        add_para(tf, text, font_size=size, bold=bold, color=color)


def slide_augmented_queries(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Escalabilidade das Consultas Aumentadas", ACCENT)

    img = GRAFICOS / "escala_consultas_aumentadas.png"
    add_image(slide, img, 0.2, 1.15, 8.8)

    tf = add_content_box(slide, 9.2, 1.3, 3.9, 5.8)
    pts = [
        ("rank · select · range_agg", True, 16, ACCENT),
        ("", False, 8, WHITE),
        ("AVL (azul):", True, 15, AVL_C),
        ("segue O(log n) — usa size e\nsubtree_sum armazenados", False, 14, WHITE),
        ("", False, 8, WHITE),
        ("BST (vermelho):", True, 15, BST_C),
        ("rank/select: O(n) — percorre\ntoda a subárvore para contar", False, 14, WHITE),
        ("range_agg BST: visita\ntodos os nós no intervalo", False, 14, WHITE),
        ("", False, 8, WHITE),
        ("Divergência clara:", True, 15, rgb(0x90, 0xEE, 0x90)),
        ("em n=1M, BST ≈ 1000× mais\nlenta que AVL no rank", False, 14, GRAY),
        ("", False, 8, WHITE),
        ("Este é o ganho real\nda augmentação!", True, 16, ACCENT),
    ]
    for text, bold, size, color in pts:
        add_para(tf, text, font_size=size, bold=bold, color=color)


def slide_theta_order(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Theta (localidade) e Ordem de Inserção", ACCENT)

    # Duas imagens lado a lado
    img_theta = GRAFICOS / "theta_latencias_avl.png"
    img_order = GRAFICOS / "ordem_avl_bst.png"
    add_image(slide, img_theta, 0.2, 1.15, 6.5)
    add_image(slide, img_order, 6.8, 1.15, 6.3)

    add_textbox(slide,
                "Theta: θ=1,2 concentra 80% dos acessos em 20% das chaves "
                "(Zipfian). Operações na AVL ficam ~30% mais rápidas "
                "por localidade de cache (chaves quentes próximas da raiz).",
                0.2, 5.6, 6.5, 1.7, font_size=13, color=GRAY)

    add_textbox(slide,
                "Ordem de inserção: BST sorted degenera (h = n = 1M+). "
                "AVL mantém h ≤ 30 em ambos os casos — "
                "rebalanceamento torna o pior caso impossível.",
                6.8, 5.6, 6.3, 1.7, font_size=13, color=GRAY)


def slide_conclusion(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, DARK)
    section_bar(slide, "Conclusão", AVL_C)

    tf = add_content_box(slide, 0.5, 1.3, 8.5, 5.8)
    items = [
        ("Contribuições do trabalho:", True, 19, ACCENT),
        ("  AVL aumentada com size e subtree_sum — O(log n) garantido", False, 17, WHITE),
        ("  rank · select · range_agg sem varredura linear", False, 17, WHITE),
        ("  Suite de testes + validação contínua de invariantes", False, 17, WHITE),
        ("  Benchmark em 11 cenários × 2 estruturas × 3 repetições", False, 17, WHITE),
        ("", False, 10, WHITE),
        ("Resultados experimentais confirmam:", True, 19, ACCENT),
        ("  Altura AVL ≤ 1,44 log₂ n em todas as escalas", False, 17, rgb(0x90, 0xEE, 0x90)),
        ("  Consultas aumentadas: AVL até 1000× mais rápida (n=1M)", False, 17, rgb(0x90, 0xEE, 0x90)),
        ("  BST sorted → colapso total (h=n)", False, 17, BST_C),
        ("  Localidade Zipfian (θ=1,2) melhora ~30% a AVL", False, 17, rgb(0x90, 0xEE, 0x90)),
        ("", False, 10, WHITE),
        ("Trade-off:", True, 19, ACCENT),
        ("  Inserção AVL ~10× mais custosa (rotações)", False, 17, GRAY),
        ("  Justificado quando consultas dominam o workload", False, 17, GRAY),
    ]
    for text, bold, size, color in items:
        add_para(tf, text, font_size=size, bold=bold, color=color)

    add_textbox(slide, "Obrigado!", 9.2, 5.0, 3.8, 1.2,
                font_size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Grupo 16 — Estruturas de Dados 2026",
                9.2, 6.1, 3.8, 0.8, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_capa(prs)
    slide_node_structure(prs)
    slide_rotations(prs)
    slide_augmented_ops(prs)
    slide_tests(prs)
    slide_transition(prs)
    slide_height(prs)
    slide_basic_ops(prs)
    slide_augmented_queries(prs)
    slide_theta_order(prs)
    slide_conclusion(prs)

    out = Path(__file__).parent / "apresentacao_grupo16.pptx"
    prs.save(str(out))
    print(f"[ok] {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
